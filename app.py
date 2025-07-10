#!/usr/bin/env python3
"""
Stock Pitch AI - Complete Version with Free and Premium Options
Full featured application with PPT generation and choice between analysis modes
"""

import streamlit as st
import sys
import os
import yfinance as yf
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import pandas as pd
from datetime import datetime
from typing import Dict, Any

# Add src directory to path
sys.path.append('src')

from data_analysis.free_analyzer import FreeStockAnalyzer
from data_analysis.stock_analyzer import StockAnalyzer
from presentation.pitch_generator_free import create_presentation_free
from presentation.pitch_generator_ai import create_presentation_ai
from utils.config import Config

def create_basic_presentation(symbol: str, analysis_results: Dict[str, Any], info: Dict[str, Any]) -> str:
    """
    Create a basic PowerPoint presentation without AI enhancement.
    
    Args:
        symbol: Stock ticker symbol
        analysis_results: Results from stock analysis
        info: Company information from yfinance
        
    Returns:
        Path to the generated presentation file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = f"{symbol} Stock Pitch"
        subtitle.text = f"Investment Analysis - {info.get('longName', symbol)}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_title = summary_slide.shapes.title
        summary_title.text = "Executive Summary"
        
        # Add content to summary slide
        content = summary_slide.placeholders[1]
        content_text = f"""
        Company: {info.get('longName', symbol)}
        Sector: {info.get('sector', 'Unknown')}
        Current Price: ${info.get('currentPrice', 0):.2f}
        Market Cap: ${info.get('marketCap', 0) / 1e9:.1f}B
        
        Recommendation: {analysis_results.get('recommendation', 'HOLD')}
        Target Price: ${analysis_results.get('target_price', 0):.2f}
        Upside Potential: {analysis_results.get('upside_potential', 'N/A')}%
        """
        content.text = content_text
        
        # Set font for content
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
        
        # Valuation slide
        valuation_slide = prs.slides.add_slide(prs.slide_layouts[1])
        val_title = valuation_slide.shapes.title
        val_title.text = "Valuation Analysis"
        
        val_content = valuation_slide.placeholders[1]
        
        # Extract valuation info from AI analysis
        highlights = analysis_results.get('highlights', [])
        target_price = analysis_results.get('target_price', 0)
        upside_potential = analysis_results.get('upside_potential', 'N/A')
        recommendation = analysis_results.get('recommendation', 'HOLD')
        
        # Build valuation content from AI analysis
        valuation_text = f"""
        TARGET PRICE: ${target_price:.2f}
        CURRENT PRICE: ${info.get('currentPrice', 0):.2f}
        UPSIDE POTENTIAL: {upside_potential}%
        RECOMMENDATION: {recommendation}
        
        KEY VALUATION METRICS:
        """
        
        # Add highlights as bullet points
        for highlight in highlights[:5]:  # Limit to 5 highlights
            valuation_text += f"• {highlight}\n"
        
        val_content.text = valuation_text
        
        # Set font for valuation content
        for paragraph in val_content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
        
        # Investment Thesis slide
        thesis_slide = prs.slides.add_slide(prs.slide_layouts[1])
        thesis_title = thesis_slide.shapes.title
        thesis_title.text = "Investment Thesis"
        
        thesis_content = thesis_slide.placeholders[1]
        thesis_content.text = analysis_results.get('investment_thesis', 'Investment based on fundamental analysis.')
        
        # Set font for thesis content
        for paragraph in thesis_content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
        
        # Key Highlights slide
        highlights_slide = prs.slides.add_slide(prs.slide_layouts[1])
        highlights_title = highlights_slide.shapes.title
        highlights_title.text = "Key Highlights & Risks"
        
        highlights_content = highlights_slide.placeholders[1]
        highlights = analysis_results.get('highlights', [])
        risks = analysis_results.get('risks', [])
        
        content_text = "KEY HIGHLIGHTS:\n"
        for highlight in highlights[:4]:  # Top 4 highlights
            content_text += f"• {highlight}\n"
        
        content_text += "\nKEY RISKS:\n"
        for risk in risks[:4]:  # Top 4 risks
            content_text += f"• {risk}\n"
        
        highlights_content.text = content_text
        
        # Set font for highlights content
        for paragraph in highlights_content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
        
        # Detailed Analysis slide
        analysis_slide = prs.slides.add_slide(prs.slide_layouts[1])
        analysis_title = analysis_slide.shapes.title
        analysis_title.text = "Detailed Financial Analysis"
        
        analysis_content = analysis_slide.placeholders[1]
        full_analysis = analysis_results.get('analysis', 'Analysis not available')
        
        # Truncate if too long for PowerPoint
        if len(full_analysis) > 1000:
            full_analysis = full_analysis[:1000] + "..."
        
        analysis_content.text = full_analysis
        
        # Set font for analysis content
        for paragraph in analysis_content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(12)
        
        # Save presentation
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        
        filename = f"{symbol}_basic_pitch_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(output_dir, filename)
        
        prs.save(filepath)
        return filepath
        
    except Exception as e:
        print(f"Error creating basic presentation: {str(e)}")
        return None

# Configure page
st.set_page_config(
    page_title="Stock Pitch AI - Professional Analysis",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS - Enhanced UI
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 15px;
        text-align: center;
        color: white;
        margin-bottom: 2rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    
    .metric-container {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        padding: 1.5rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 1rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        transition: transform 0.3s ease;
    }
    
    .metric-container:hover {
        transform: translateY(-5px);
    }
    
    .analysis-section {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 15px;
        margin-bottom: 2rem;
        color: white;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    
    .recommendation-buy {
        background: linear-gradient(135deg, #4CAF50 0%, #45a049 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin: 1rem 0;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    
    .recommendation-sell {
        background: linear-gradient(135deg, #f44336 0%, #da190b 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin: 1rem 0;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    
    .recommendation-hold {
        background: linear-gradient(135deg, #ff9800 0%, #f57c00 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin: 1rem 0;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    
    .feature-card {
        background: #e8ebf0;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 5px solid #667eea;
        margin-bottom: 1rem;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        color: #2c3e50;
        font-weight: 500;
    }
    
    .ppt-generation-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 15px;
        margin-bottom: 2rem;
        color: white;
        box-shadow: 0 6px 12px rgba(0, 0, 0, 0.2);
        border: 3px solid #4facfe;
    }
    
    .ppt-generation-card h3 {
        color: white;
        margin-bottom: 1rem;
    }
    
    .ppt-generation-card p {
        color: #f0f8ff;
        margin-bottom: 0.5rem;
    }
    
    .sidebar-section {
        background: #e8ebf0;
        padding: 1rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        color: #2c3e50;
        font-weight: 500;
    }
    
    .dcf-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.5rem;
        border-radius: 10px;
        color: white;
        margin-bottom: 1rem;
    }
    
    .wacc-card {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        padding: 1.5rem;
        border-radius: 10px;
        color: white;
        margin-bottom: 1rem;
    }
    
    .financial-card {
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        padding: 1.5rem;
        border-radius: 10px;
        color: white;
        margin-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)

def main():
    # Header
    st.markdown("""
    <div class="main-header">
        <h1>🚀 Stock Pitch AI - Professional Analysis Platform</h1>
        <p>Advanced Financial Analysis with DCF, WACC & Financial Statement Analysis</p>
        <p>Complete with PowerPoint Generation & AI-Powered Insights</p>
    </div>
    """, unsafe_allow_html=True)

    # Initialize session state
    if 'analysis_complete' not in st.session_state:
        st.session_state.analysis_complete = False
    if 'analysis_result' not in st.session_state:
        st.session_state.analysis_result = None
    if 'stock_info' not in st.session_state:
        st.session_state.stock_info = None
    if 'current_symbol' not in st.session_state:
        st.session_state.current_symbol = None

    # Sidebar - Enhanced
    with st.sidebar:
        st.markdown("""
        <div class="sidebar-section">
            <h2>🎯 Stock Pitch AI Configuration</h2>
            <p>Configure your analysis parameters below</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Analysis Mode Selection
        st.markdown("""
        <div class="sidebar-section">
            <h3>� Analysis Mode</h3>
            <p>Choose between free comprehensive analysis or premium AI-powered insights</p>
        </div>
        """, unsafe_allow_html=True)
        
        analysis_mode = st.selectbox(
            "Select Analysis Mode",
            options=["Free Analysis (No API Required)", "Premium Analysis (with OpenAI GPT)"],
            index=0,
            help="Free mode: Advanced financial analysis using built-in algorithms\nPremium mode: AI-enhanced insights with OpenAI GPT integration"
        )
        
        # API Key input for premium mode
        api_key = None
        if "Premium" in analysis_mode:
            st.markdown("""
            <div class="sidebar-section">
                <h4>🔑 OpenAI API Key Required</h4>
                <p>Enter your OpenAI API key to unlock AI-powered insights</p>
            </div>
            """, unsafe_allow_html=True)
            
            api_key = st.text_input(
                "OpenAI API Key",
                type="password",
                placeholder="sk-...",
                help="Get your API key from https://platform.openai.com/api-keys"
            )
            
            if not api_key:
                st.warning("⚠️ Please enter your OpenAI API key to use premium features")
                st.info("💡 No API key? Switch to Free Analysis for comprehensive analysis without API requirements")
        
        # Stock Selection
        st.markdown("""
        <div class="sidebar-section">
            <h3>📈 Stock Selection</h3>
            <p>Enter the stock ticker symbol you want to analyze</p>
        </div>
        """, unsafe_allow_html=True)
        
        stock_symbol = st.text_input(
            "Stock Symbol",
            placeholder="Enter symbol (e.g., AAPL, MSFT, GOOGL)",
            help="Enter the stock ticker symbol (e.g., AAPL for Apple, MSFT for Microsoft)",
            value=""
        ).upper()
        
        # Clear session state when switching stocks
        if st.session_state.current_symbol != stock_symbol:
            st.session_state.analysis_complete = False
            st.session_state.analysis_result = None
            st.session_state.stock_info = None
            st.session_state.current_symbol = stock_symbol
        
        # Analysis Parameters
        st.markdown("""
        <div class="sidebar-section">
            <h3>⚙️ Advanced Analysis Parameters</h3>
            <p>Customize your financial analysis settings</p>
        </div>
        """, unsafe_allow_html=True)
        
        dcf_years = st.slider(
            "DCF Projection Years",
            min_value=3,
            max_value=10,
            value=5,
            help="📊 DCF (Discounted Cash Flow) Projects future cash flows for this many years to calculate intrinsic value. More years = longer-term view."
        )
        
        st.markdown("**📈 What is DCF?** DCF estimates a stock's fair value by projecting future cash flows and discounting them to present value.")
        
        risk_free_rate = st.slider(
            "Risk-Free Rate (%)",
            min_value=2.0,
            max_value=6.0,
            value=4.5,
            step=0.1,
            help="🏛️ Risk-Free Rate: The return on a 'risk-free' investment (like US Treasury bonds). Used as baseline for calculating required returns."
        )
        
        st.markdown("**🏛️ What is Risk-Free Rate?** The theoretical return on an investment with zero risk, typically based on 10-year US Treasury bonds.")
        
        market_risk_premium = st.slider(
            "Market Risk Premium (%)",
            min_value=4.0,
            max_value=8.0,
            value=6.5,
            step=0.1,
            help="📊 Market Risk Premium: Extra return investors demand for taking market risk vs. risk-free investments. Higher = more conservative."
        )
        
        st.markdown("**📊 What is Market Risk Premium?** The additional return investors expect for taking on market risk instead of risk-free investments.")
        
        # Features showcase
        st.markdown("---")
        st.markdown("""
        <div class="sidebar-section">
            <h3>🎯 Analysis Features</h3>
            <p>What you'll get with your analysis:</p>
            <ul style="margin-left: 1rem;">
                <li>✅ DCF Valuation (Fair Value Calculation)</li>
                <li>✅ WACC Calculation (Cost of Capital)</li>
                <li>✅ Financial Health Assessment</li>
                <li>✅ Risk Analysis & Beta Calculation</li>
                <li>✅ Investment Recommendation</li>
                <li>✅ Professional PowerPoint Report</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)

    # Main content
    if stock_symbol:
        try:
            # Fetch stock data
            with st.spinner(f"📊 Fetching comprehensive data for {stock_symbol}..."):
                ticker = yf.Ticker(stock_symbol)
                info = ticker.info
                hist = ticker.history(period="2y")
                financials = ticker.financials if hasattr(ticker, 'financials') else pd.DataFrame()
                
            # Display enhanced stock info
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.markdown(f"""
                <div class="metric-container">
                    <h3>{info.get('longName', stock_symbol)}</h3>
                    <p><strong>{stock_symbol}</strong></p>
                    <p>{info.get('sector', 'Unknown Sector')}</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                current_price = info.get('currentPrice', 0)
                prev_close = info.get('previousClose', current_price)
                change = ((current_price - prev_close) / prev_close * 100) if prev_close else 0
                color = "green" if change >= 0 else "red"
                arrow = "↗" if change >= 0 else "↘"
                
                st.markdown(f"""
                <div class="metric-container">
                    <h4>Current Price</h4>
                    <h2>${current_price:.2f}</h2>
                    <p style="color: {color};">{arrow} {change:+.2f}%</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                market_cap = info.get('marketCap', 0)
                market_cap_b = market_cap / 1e9 if market_cap > 0 else 0
                cap_type = "Large Cap" if market_cap_b > 10 else "Mid Cap" if market_cap_b > 2 else "Small Cap"
                
                st.markdown(f"""
                <div class="metric-container">
                    <h4>Market Cap</h4>
                    <h2>${market_cap_b:.1f}B</h2>
                    <p>{cap_type}</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col4:
                pe_ratio = info.get('trailingPE', 0)
                pe_assessment = "Low" if pe_ratio < 15 else "High" if pe_ratio > 30 else "Moderate"
                
                st.markdown(f"""
                <div class="metric-container">
                    <h4>P/E Ratio</h4>
                    <h2>{pe_ratio:.1f}x</h2>
                    <p>{pe_assessment} Valuation</p>
                </div>
                """, unsafe_allow_html=True)
            
            # Additional metrics row
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                beta = info.get('beta', 1.0)
                volatility = "Low" if beta < 1 else "High" if beta > 1.5 else "Moderate"
                st.metric("Beta (Risk)", f"{beta:.2f}", f"{volatility} Volatility")
            
            with col2:
                dividend_yield = info.get('dividendYield', 0)
                div_pct = dividend_yield * 100 if dividend_yield else 0
                st.metric("Dividend Yield", f"{div_pct:.2f}%")
            
            with col3:
                high_52w = info.get('fiftyTwoWeekHigh', 0)
                low_52w = info.get('fiftyTwoWeekLow', 0)
                st.metric("52W High", f"${high_52w:.2f}")
            
            with col4:
                st.metric("52W Low", f"${low_52w:.2f}")
                
            # Enhanced price chart with volume
            st.markdown("### 📈 Stock Performance & Technical Analysis")
            
            # Create subplots
            fig = make_subplots(
                rows=2, cols=1,
                shared_xaxes=True,
                vertical_spacing=0.1,
                row_width=[0.7, 0.3],
                subplot_titles=('Stock Price', 'Volume')
            )
            
            # Price chart
            fig.add_trace(
                go.Scatter(
                    x=hist.index,
                    y=hist['Close'],
                    mode='lines',
                    name=f'{stock_symbol} Price',
                    line=dict(color='#667eea', width=3),
                    hovertemplate='<b>%{y:.2f}</b><br>%{x}<extra></extra>'
                ),
                row=1, col=1
            )
            
            # Volume chart
            fig.add_trace(
                go.Bar(
                    x=hist.index,
                    y=hist['Volume'],
                    name='Volume',
                    marker_color='rgba(102, 126, 234, 0.6)',
                    hovertemplate='<b>%{y:,.0f}</b><br>%{x}<extra></extra>'
                ),
                row=2, col=1
            )
            
            fig.update_layout(
                title=f"{stock_symbol} - 2 Year Performance Analysis",
                xaxis_title="Date",
                height=600,
                showlegend=False
            )
            
            fig.update_yaxes(title_text="Price ($)", row=1, col=1)
            fig.update_yaxes(title_text="Volume", row=2, col=1)
            
            st.plotly_chart(fig, use_container_width=True)
            
            # Analysis button
            if st.button("🔍 Generate Comprehensive Analysis", key="analyze"):
                with st.spinner("🔬 Running advanced financial analysis..."):
                    # Prepare stock data
                    stock_data = {
                        'symbol': stock_symbol,
                        'company_name': info.get('longName', stock_symbol),
                        'current_price': info.get('currentPrice', 0),
                        'pe_ratio': info.get('trailingPE', 0),
                        'eps': info.get('trailingEps', 0),
                        'market_cap': info.get('marketCap', 0),
                        'beta': info.get('beta', 1.0),
                        'dividend_yield': info.get('dividendYield', 0),
                        'pb_ratio': info.get('priceToBook', 0),
                        '52w_high': info.get('fiftyTwoWeekHigh', 0),
                        '52w_low': info.get('fiftyTwoWeekLow', 0),
                        'sector': info.get('sector', 'Unknown')
                    }
                    
                    # Run analysis based on selected mode
                    if "Free Analysis" in analysis_mode:
                        analyzer = FreeStockAnalyzer()
                        result = analyzer.analyze_stock_free(stock_data)
                        analysis_type = "Free"
                    else:  # Premium Analysis
                        if api_key:
                            try:
                                config = Config(api_key=api_key)
                                analyzer = StockAnalyzer(config)
                                # First fetch stock data, then analyze
                                premium_stock_data = analyzer.fetch_stock_data(stock_symbol)
                                result = analyzer.analyze_stock(premium_stock_data)
                                analysis_type = "Premium"
                            except Exception as e:
                                st.error(f"❌ Premium analysis failed: {str(e)}")
                                st.info("🔄 Falling back to free analysis...")
                                analyzer = FreeStockAnalyzer()
                                result = analyzer.analyze_stock_free(stock_data)
                                analysis_type = "Free (Fallback)"
                        else:
                            st.error("❌ API key required for premium analysis")
                            return
                    
                    # Store results in session state
                    st.session_state.analysis_complete = True
                    st.session_state.analysis_result = result
                    st.session_state.stock_info = info
                    st.session_state.current_symbol = stock_symbol
                    st.session_state.analysis_type = analysis_type  # <-- Store analysis_type

                    # Display results
                    st.success(f"✅ {analysis_type} Analysis Complete!")
                    
                    # Key metrics
                    if analysis_type == "Premium":
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            recommendation = result.get('recommendation', 'N/A')
                            st.markdown(f"""
                            <div style="background: #4CAF50; color: white; padding: 1rem; border-radius: 5px; text-align: center;">
                                <h3>Recommendation</h3>
                                <h2>{recommendation}</h2>
                            </div>
                            """, unsafe_allow_html=True)
                        with col2:
                            target_price = result.get('target_price', 'N/A')
                            st.markdown(f"""
                            <div style="background: #667eea; color: white; padding: 1rem; border-radius: 5px; text-align: center;">
                                <h3>Target Price</h3>
                                <h2>${target_price if target_price != 'N/A' else 'N/A'}</h2>
                            </div>
                            """, unsafe_allow_html=True)
                        with col3:
                            upside = result.get('upside_potential', 'N/A')
                            st.markdown(f"""
                            <div style="background: #764ba2; color: white; padding: 1rem; border-radius: 5px; text-align: center;">
                                <h3>Upside Potential</h3>
                                <h2>{upside if upside != 'N/A' else 'N/A'}</h2>
                            </div>
                            """, unsafe_allow_html=True)

                    # DCF and WACC Analysis (AI version)
                    if analysis_type == "Premium":
                        st.markdown("## 💰 DCF & WACC Analysis (AI)")
                        ai_analysis = result.get('analysis', None)
                        if ai_analysis:
                            st.markdown(f"""
                            <div class="analysis-section">
                                <pre style='font-family:Consolas,monospace;font-size:15px;background:#222;color:#fff;padding:1em;border-radius:8px;'>{ai_analysis}</pre>
                            </div>
                            """, unsafe_allow_html=True)
                        else:
                            st.warning("No DCF/WACC analysis returned by AI.")

                    # Investment Thesis (AI version)
                    if analysis_type == "Premium":
                        st.markdown("## 🎯 Investment Thesis (AI)")
                        thesis = result.get('investment_thesis', None)
                        if thesis:
                            st.markdown(f"<div class='analysis-section'><p><strong>{thesis}</strong></p></div>", unsafe_allow_html=True)
                        else:
                            st.warning("No investment thesis returned by AI.")

                    # Key Highlights and Risks (AI version)
                    if analysis_type == "Premium":
                        col1, col2 = st.columns(2)
                        with col1:
                            st.markdown("### ✅ Key Highlights (AI)")
                            highlights = result.get('highlights', [])
                            if highlights:
                                for highlight in highlights:
                                    st.write(f"• {highlight}")
                            else:
                                st.warning("No highlights returned by AI.")
                        with col2:
                            st.markdown("### ⚠️ Key Risks (AI)")
                            risks = result.get('risks', [])
                            if risks:
                                for risk in risks:
                                    st.write(f"• {risk}")
                            else:
                                st.warning("No risks returned by AI.")
                    
        except Exception as e:
            st.error(f"❌ Error analyzing {stock_symbol}: {str(e)}")
    
    # PowerPoint Generation Section (available after analysis)
    if st.session_state.analysis_complete and st.session_state.analysis_result:
        st.markdown("---")
        st.markdown("## 📊 Generate Professional PowerPoint Presentation")
        
        # Add option to start new analysis
        col_clear, col_spacer = st.columns([1, 4])
        with col_clear:
            if st.button("🔄 New Analysis", help="Clear current analysis and start fresh"):
                st.session_state.analysis_complete = False
                st.session_state.analysis_result = None
                st.session_state.stock_info = None
                st.session_state.current_symbol = None
                st.rerun()
        
        st.markdown("""
        <div class="ppt-generation-card">
            <h3>🎯 Professional Stock Pitch Presentation</h3>
            <p><strong>Generate a comprehensive PowerPoint presentation</strong> with all analysis results, charts, and recommendations ready for professional use.</p>
            <p>✅ Includes all analysis data | ✅ Professional formatting | ✅ Ready to download</p>
        </div>
        """, unsafe_allow_html=True)
        
        col1, col2 = st.columns([1, 1])
        
        with col1:
            presentation_style = st.selectbox(
                "📋 Presentation Style",
                options=["Professional", "Detailed", "Executive Summary"],
                index=0,
                help="Choose the style and depth of the presentation",
                key="ppt_style"
            )
        
        with col2:
            st.markdown("### 🎨 Generate Your Presentation")
            
        # Large, prominent button
        if st.button("🚀 GENERATE POWERPOINT PRESENTATION", 
                    type="primary", 
                    use_container_width=True,
                    help="Click to generate a professional PowerPoint presentation with all your analysis",
                    key="generate_ppt"):
            
            with st.spinner("🎨 Creating your professional presentation..."):
                try:
                    # Get data from session state
                    result = st.session_state.analysis_result
                    info = st.session_state.stock_info
                    symbol = st.session_state.current_symbol
                    
                    # Use the proper PitchGenerator class instead of basic function
                    st.info("🔄 Creating presentation...")
                    analysis_type = st.session_state.get('analysis_type', 'Free')
                    if analysis_type == "Premium":
                        ppt_path = create_presentation_ai(stock_symbol, result)
                    else:
                        ppt_path = create_presentation_free(stock_symbol, result, info)
                    
                    if ppt_path and os.path.exists(ppt_path):
                        st.success("✅ PowerPoint presentation created successfully!")
                        
                        # Provide download link
                        with open(ppt_path, "rb") as file:
                            st.download_button(
                                label="📥 Download PowerPoint Presentation",
                                data=file.read(),
                                file_name=f"{symbol}_stock_pitch_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                type="primary",
                                use_container_width=True,
                                key="download_ppt"
                            )
                        
                        # Show presentation details
                        file_size = os.path.getsize(ppt_path)
                        st.markdown(f"""
                        <div class="ppt-generation-card">
                            <h4>📊 Presentation Details</h4>
                            <p><strong>File:</strong> {os.path.basename(ppt_path)}</p>
                            <p><strong>Size:</strong> {file_size:,} bytes</p>
                            <p><strong>Format:</strong> PowerPoint (.pptx)</p>
                            <p><strong>Content:</strong> Complete stock analysis with charts and recommendations</p>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        st.balloons()  # Celebration animation
                    else:
                        st.error("❌ Unable to create presentation. Please try again.")
                        
                except Exception as e:
                    st.error(f"❌ Error generating presentation: {str(e)}")
    
    # Welcome section when no analysis is done
    elif not stock_symbol:
        st.markdown("## 👋 Welcome to Stock Pitch AI")
        st.markdown("""
        Enter a stock symbol in the sidebar to get started with comprehensive financial analysis including:
        
        - **DCF Valuation** - Discounted Cash Flow modeling
        - **WACC Calculation** - Weighted Average Cost of Capital
        - **Financial Statement Analysis** - Comprehensive health assessment
        - **Comparative Valuation** - Industry benchmarking
        - **Risk Assessment** - Investment risk evaluation
        
        **No API keys required!** This free version provides professional-grade analysis.
        """)

if __name__ == "__main__":
    main()
