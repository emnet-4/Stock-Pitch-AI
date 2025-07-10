"""
AI Version Presentation Generator
Generates a PowerPoint using AI-powered premium analysis results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from typing import Dict, Any
import os

def create_presentation_ai(symbol: str, analysis_results: Dict[str, Any]) -> str:
    prs = Presentation()
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = f"{symbol} Stock Pitch (AI Mode)"
    subtitle.text = f"AI-Powered Analysis\n{datetime.now().strftime('%B %d, %Y')}"
    # Executive Summary
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_title = summary_slide.shapes.title
    summary_title.text = "Executive Summary"
    content = summary_slide.placeholders[1]
    tf = content.text_frame
    p = tf.paragraphs[0]
    thesis = analysis_results.get('investment_thesis', 'No investment thesis returned by AI.')
    p.text = thesis
    p.font.size = Pt(16)
    # Main AI Analysis
    analysis = analysis_results.get('analysis', None)
    if analysis:
        detail_slide = prs.slides.add_slide(prs.slide_layouts[1])
        detail_slide.shapes.title.text = "DCF & WACC Analysis (AI)"
        content = detail_slide.placeholders[1]
        tf = content.text_frame
        p = tf.paragraphs[0]
        p.text = analysis
        p.font.size = Pt(13)
        p.font.name = 'Consolas'
    # Highlights & Risks Slide
    highlights = analysis_results.get('highlights', [])
    risks = analysis_results.get('risks', [])
    if highlights or risks:
        hr_slide = prs.slides.add_slide(prs.slide_layouts[1])
        hr_slide.shapes.title.text = "Highlights & Risks (AI)"
        content = hr_slide.placeholders[1]
        tf = content.text_frame
        if highlights:
            p = tf.paragraphs[0]
            p.text = "Key Highlights:"
            p.font.size = Pt(15)
            p.font.bold = True
            for h in highlights:
                p = tf.add_paragraph()
                p.text = f"• {h}"
                p.font.size = Pt(13)
                p.level = 1
        if risks:
            p = tf.add_paragraph()
            p.text = "Key Risks:"
            p.font.size = Pt(15)
            p.font.bold = True
            for r in risks:
                p = tf.add_paragraph()
                p.text = f"• {r}"
                p.font.size = Pt(13)
                p.level = 1
    # Recommendation Slide
    rec_slide = prs.slides.add_slide(prs.slide_layouts[1])
    rec_slide.shapes.title.text = "Investment Recommendation (AI)"
    content = rec_slide.placeholders[1]
    tf = content.text_frame
    rec = analysis_results.get('recommendation', 'N/A')
    p = tf.paragraphs[0]
    p.text = f"Recommendation: {rec}"
    p.font.size = Pt(22)
    p.font.bold = True
    target_price = analysis_results.get('target_price', 'N/A')
    upside = analysis_results.get('upside_potential', 'N/A')
    p = tf.add_paragraph()
    p.text = f"Target Price: {target_price}"
    p.font.size = Pt(15)
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Upside Potential: {upside}"
    p.font.size = Pt(15)
    p.level = 1
    thesis = analysis_results.get('investment_thesis', None)
    if thesis:
        p = tf.add_paragraph()
        p.text = thesis
        p.font.size = Pt(14)
        p.level = 1
    # Save
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{symbol}_stock_pitch_ai_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath
